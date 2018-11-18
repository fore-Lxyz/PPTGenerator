#
from pptx import Presentation
from pptx.util import Inches, Pt
from werkzeug.utils import secure_filename
import os

#filePath = '/Users/luomei/Desktop/test.md'

def autoCreatPPT(filePath):
    endLine = '---'
    with open(filePath) as file:
        lines = file.readlines()
    lines = [i.strip().rstrip('\n') for i in lines]
    #print(lines)
    data = []
    titles = {}
    contents = []
    levels = []
    paras = []
    prs = Presentation()

    while '' in lines:
        lines.remove('')
    line = lines[len(lines) - 1]
    #print(line)
    if line != endLine:
        lines.append('---')
    for line in lines:
        level = 0
        if line == '':
            continue
        #print('line', line)
        if line == endLine:
            curLevel = 0
            curHolders = 1
            if len(levels) == 1:
                if levels[0] == 1:
                    title_slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    title.text = contents[0]
                else:
                    title_slide_layout = prs.slide_layouts[2]
                    slide = prs.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    title.text = contents[0]
            elif len(levels) == 2:
                if levels[1] > 0:
                    title_slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    title.text = contents[0]
                    subtitle = slide.placeholders[1]
                    subtitle.text = contents[1]
                elif levels[0] > 0:
                    title_slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    title.text = contents[0]
                    subtitle = slide.shapes.placeholders[1]
                    subtitle.text_frame.text = contents[1]
                else:
                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)
                    left = top = width = height = Inches(2)
                    textBox = slide.shapes.add_textbox(left, top, width, height)
                    textBox.text_frame.text = contents[0]
                    para = textBox.text_frame.add_paragraph()
                    para.text = contents[1]

            else:
                title_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                for i in range(len(levels)):
                    #print(i, levels[i], contents[i])
                    if levels[i] == 1:
                        title.text = contents[i]
                    elif levels[i] == 2:
                        subtitle = slide.shapes.placeholders[1]
                        if subtitle.text_frame.text != '':
                            para = subtitle.text_frame.add_paragraph()
                            para.text = contents[i]
                            para.level = 0
                        else:
                            subtitle.text_frame.text = contents[i]
                    elif levels[i] < 0:
                        para = subtitle.text_frame.add_paragraph()
                        para.text = contents[i]
                        para.level = abs(levels[i]) + curLevel
                    elif levels[i] > 2:
                        para = subtitle.text_frame.add_paragraph()
                        para.text = contents[i]
                        para.level = abs(levels[i]-2)
                        curLevel = levels[i] - 2
            contents.clear()
            levels.clear()
            titles.clear()
            paras.clear()
            continue
        for ch in line:
            if ch == '#':
                level += 1
            else:
                if level != 0:
                    break
        if level != 0:
            line = line.strip().lstrip('#')
            #print(line)
            titles[line] = level
            contents.append(line)
            levels.append(level)
        else:
            paras.append(line)
            contents.append(line)
            levels.append(-1)
            # title_slide_layout = prs.slide_layouts[level-1]
            # slide = prs.slides.add_slide(title_slide_layout)
            # title = slide.shapes.title
            # title.text = line

    prs.save(os.path.join(os.path.dirname(__file__), 'static/ppt',
                                   secure_filename("test.pptx")))
    print("success to save the ppt")


